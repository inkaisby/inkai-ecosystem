#!/usr/bin/env python3
"""Generate INKAI Mobile Web proposal PowerPoint — presidential edition."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from capture_ppt_screenshots import asset_path, capture_all  # noqa: E402

# Presidential palette — navy, champagne gold, cream
NAVY = RGBColor(0x0B, 0x1A, 0x2E)
NAVY_MID = RGBColor(0x14, 0x28, 0x42)
PANEL = RGBColor(0x10, 0x22, 0x38)
PANEL_ALT = RGBColor(0x0E, 0x1E, 0x34)
GOLD = RGBColor(0xC9, 0xA9, 0x62)
GOLD_DIM = RGBColor(0xA8, 0x88, 0x48)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
SILVER = RGBColor(0xB8, 0xC4, 0xCE)
MUTED = RGBColor(0x8A, 0x96, 0xA8)
WHITE = CREAM
GRAY = SILVER
GRAY_DARK = MUTED
BG_DARK = NAVY
BG_CARD = PANEL_ALT

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "INKAI-Mobile-Web-Proposal.pptx"
LOGO_PATH = Path(r"D:\website\inkai\logo.png")
PHOTO_W = Inches(5.1)
CONTENT_X = Inches(5.28)
CONTENT_W = Inches(7.75)
FOOTER_TEXT = "Institut Karate-Do Indonesia  |  INKAI Digital Ecosystem  |  inkai-mobile-web.vercel.app"

# Typography & branding — line spacing 2.5, logo besar proporsional
LINE_SPACING = 2.5
LOGO_TITLE_IN = 1.85      # slide judul / penutup
LOGO_SECTION_IN = 1.15      # slide pembagi bagian
LOGO_CONTENT_IN = 0.72      # slide konten (pojok kanan)
TABLE_ROW_MIN_IN = 0.48     # tinggi baris tabel (sesuai line spacing 2.5)


def _para(p, size: int, color: RGBColor, *, bold: bool = False, align=None, spacing: float = LINE_SPACING) -> None:
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.line_spacing = spacing
    p.space_before = Pt(0)
    p.space_after = Pt(max(3, int(size * 0.2)))
    if align is not None:
        p.alignment = align


def _style_cell_paragraphs(cell, size: int, color: RGBColor, *, bold: bool = False, align=None) -> None:
    for p in cell.text_frame.paragraphs:
        _para(p, size, color, bold=bold, align=align)


def _style_textframe(tf, size: int, color: RGBColor, *, bold: bool = False) -> None:
    for p in tf.paragraphs:
        _para(p, size, color, bold=bold)


def _theme_image(theme: str) -> Path | None:
    p = asset_path(theme)
    return p if p.is_file() else (asset_path("hero") if asset_path("hero").is_file() else None)


def _rect(slide, left, top, width, height, color: RGBColor, transparency: float = 0.0):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if transparency:
        sh.fill.transparency = transparency
    sh.line.fill.background()
    return sh


def _add_photo_panel(slide, prs, theme: str) -> None:
    path = _theme_image(theme)
    if path:
        slide.shapes.add_picture(str(path), 0, 0, width=PHOTO_W, height=prs.slide_height)
        _rect(slide, PHOTO_W - Inches(0.05), 0, Inches(0.05), prs.slide_height, GOLD)
        _rect(slide, PHOTO_W - Inches(0.85), 0, Inches(0.85), prs.slide_height, NAVY, 0.5)


def _setup_split(slide, prs, theme: str) -> None:
    _add_photo_panel(slide, prs, theme)
    _rect(slide, PHOTO_W, 0, prs.slide_width - PHOTO_W, prs.slide_height, PANEL)
    _rect(slide, PHOTO_W, 0, prs.slide_width - PHOTO_W, Inches(0.06), GOLD)
    _rect(slide, PHOTO_W, prs.slide_height - Inches(0.36), prs.slide_width - PHOTO_W, Inches(0.36), NAVY_MID)


def _setup_hero(slide, prs, theme: str, overlay: float = 0.62) -> None:
    path = _theme_image(theme)
    if path:
        slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    _rect(slide, 0, 0, prs.slide_width, prs.slide_height, NAVY, overlay)
    _rect(slide, 0, 0, prs.slide_width, Inches(0.08), GOLD)
    _rect(slide, 0, prs.slide_height - Inches(0.06), prs.slide_width, Inches(0.06), GOLD_DIM)


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide, prs) -> None:
    _rect(slide, PHOTO_W if PHOTO_W else 0, 0, prs.slide_width - PHOTO_W, Inches(0.06), GOLD)


def add_footer(slide, prs, text: str = FOOTER_TEXT) -> None:
    y = prs.slide_height - Inches(0.3)
    box = slide.shapes.add_textbox(Inches(0.45), y, prs.slide_width - Inches(0.9), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    _para(p, 8, MUTED, align=PP_ALIGN.CENTER)


def add_logo(
    slide,
    prs,
    *,
    size: float | None = None,
    position: str = "top-right",
    top: float = 0.15,
) -> None:
    if not LOGO_PATH.is_file():
        return
    if size is None:
        size = LOGO_CONTENT_IN
    sz = Inches(size)
    top_in = Inches(top)
    margin = Inches(0.28)
    if position == "top-right":
        left = prs.slide_width - margin - sz
    elif position == "top-left":
        left = CONTENT_X + Inches(0.05)
    elif position == "center-top":
        left = (prs.slide_width - sz) / 2
    elif position == "section":
        left = prs.slide_width - Inches(1.05) - sz
        top_in = Inches(2.35)
    else:
        left = CONTENT_X
    slide.shapes.add_picture(str(LOGO_PATH), left, top_in, width=sz, height=sz)


def add_title_slide(prs, title: str, subtitle: str, extra: str = "", theme: str = "hero") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _setup_hero(slide, prs, theme, overlay=0.55)
    add_logo(slide, prs, size=LOGO_TITLE_IN, position="center-top", top=0.52)

    _rect(slide, Inches(5.8), Inches(2.55), Inches(1.8), Inches(0.04), GOLD)

    tbox = slide.shapes.add_textbox(Inches(1.2), Inches(2.75), Inches(11), Inches(1.2))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    _para(tp, 44, CREAM, bold=True, align=PP_ALIGN.CENTER)

    sbox = slide.shapes.add_textbox(Inches(1.2), Inches(3.85), Inches(11), Inches(0.85))
    sp = sbox.text_frame.paragraphs[0]
    sp.text = subtitle
    _para(sp, 22, GOLD, align=PP_ALIGN.CENTER)

    if extra:
        ebox = slide.shapes.add_textbox(Inches(2.0), Inches(4.75), Inches(9.3), Inches(1.6))
        ep = ebox.text_frame.paragraphs[0]
        ep.text = extra
        _para(ep, 14, SILVER, align=PP_ALIGN.CENTER)

    add_footer(slide, prs)


def add_section_slide(prs, section_num: str, title: str, subtitle: str = "", theme: str = "leadership") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _setup_hero(slide, prs, theme, overlay=0.68)
    add_logo(slide, prs, size=LOGO_SECTION_IN, position="section")

    num_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.55), Inches(3), Inches(0.55))
    np = num_box.text_frame.paragraphs[0]
    np.text = section_num
    _para(np, 16, GOLD, bold=True)

    tbox = slide.shapes.add_textbox(Inches(0.9), Inches(3.15), Inches(11.5), Inches(1.4))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    _para(tp, 40, CREAM, bold=True)

    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.9), Inches(4.45), Inches(11), Inches(0.85))
        sp = sbox.text_frame.paragraphs[0]
        sp.text = subtitle
        _para(sp, 17, SILVER)

    add_footer(slide, prs)


def add_bullet_slide(prs, title: str, bullets: list[str], subtitle: str = "", theme: str = "hero") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _setup_split(slide, prs, theme)
    add_logo(slide, prs, size=LOGO_CONTENT_IN, position="top-right", top=0.12)

    tbox = slide.shapes.add_textbox(CONTENT_X + Inches(0.12), Inches(0.28), CONTENT_W - Inches(0.2), Inches(0.7))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    _para(tp, 24, CREAM, bold=True)

    body_top = Inches(1.02)
    if subtitle:
        sbox = slide.shapes.add_textbox(CONTENT_X + Inches(0.12), Inches(0.92), CONTENT_W, Inches(0.42))
        sp = sbox.text_frame.paragraphs[0]
        sp.text = subtitle
        _para(sp, 12, GOLD)
        body_top = Inches(1.38)

    bbox = slide.shapes.add_textbox(CONTENT_X + Inches(0.18), body_top, CONTENT_W - Inches(0.25), Inches(5.55))
    tf = bbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {item}" if item.startswith("  ") else f"▸  {item}"
        sz = 13 if item.startswith("  ") else 14
        _para(p, sz, MUTED if item.startswith("  ") else SILVER)

    add_footer(slide, prs)


def add_table_slide(
    prs, title: str, headers: list[str], rows: list[list[str]], subtitle: str = "", theme: str = "executive"
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _setup_split(slide, prs, theme)
    add_logo(slide, prs, size=LOGO_CONTENT_IN, position="top-right", top=0.12)

    tbox = slide.shapes.add_textbox(CONTENT_X + Inches(0.12), Inches(0.28), CONTENT_W, Inches(0.65))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    _para(tp, 22, CREAM, bold=True)

    top = Inches(1.02)
    if subtitle:
        sbox = slide.shapes.add_textbox(CONTENT_X + Inches(0.12), Inches(0.88), CONTENT_W, Inches(0.38))
        sp = sbox.text_frame.paragraphs[0]
        sp.text = subtitle
        _para(sp, 11, GOLD)
        top = Inches(1.32)

    cols = len(headers)
    table_rows = 1 + len(rows)
    tw = CONTENT_W - Inches(0.15)
    row_h = max(Inches(TABLE_ROW_MIN_IN), Inches(5.35 / max(table_rows, 1)))
    table = slide.shapes.add_table(table_rows, cols, CONTENT_X + Inches(0.1), top, tw, row_h * table_rows).table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = GOLD_DIM
        _style_cell_paragraphs(cell, 10, CREAM, bold=True, align=PP_ALIGN.CENTER)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL_ALT if r % 2 else NAVY_MID
            color = CREAM if c == 0 else SILVER
            _style_cell_paragraphs(cell, 10, color)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_footer(slide, prs)


def add_two_column_slide(
    prs,
    title: str,
    left_title: str,
    left_items: list[str],
    right_title: str,
    right_items: list[str],
    theme: str = "stakeholder",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _setup_split(slide, prs, theme)
    add_logo(slide, prs, size=LOGO_CONTENT_IN, position="top-right", top=0.12)

    tbox = slide.shapes.add_textbox(CONTENT_X + Inches(0.12), Inches(0.28), CONTENT_W, Inches(0.65))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    _para(tp, 22, CREAM, bold=True)

    col_w = (CONTENT_W - Inches(0.35)) / 2
    for idx, (ctitle, citems) in enumerate([(left_title, left_items), (right_title, right_items)]):
        x = CONTENT_X + Inches(0.12) + idx * (col_w + Inches(0.2))
        hbox = slide.shapes.add_textbox(x, Inches(1.0), col_w, Inches(0.45))
        hp = hbox.text_frame.paragraphs[0]
        hp.text = ctitle
        _para(hp, 14, GOLD, bold=True)
        _rect(slide, x, Inches(1.45), col_w, Inches(0.02), GOLD_DIM)

        bbox = slide.shapes.add_textbox(x, Inches(1.55), col_w, Inches(5.15))
        tf = bbox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(citems):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"▸ {item}"
            _para(p, 13, SILVER)

    add_footer(slide, prs)


def fmt_idr(amount: int) -> str:
    """Format rupiah with Indonesian thousand separator."""
    return f"Rp {amount:,}".replace(",", ".")


def add_cost_table_slide(
    prs,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    subtitle: str = "",
    col_widths: list[float] | None = None,
    theme: str = "cost",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _setup_split(slide, prs, theme)
    add_logo(slide, prs, size=LOGO_CONTENT_IN, position="top-right", top=0.12)

    tbox = slide.shapes.add_textbox(CONTENT_X + Inches(0.12), Inches(0.28), CONTENT_W, Inches(0.6))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    _para(tp, 20, CREAM, bold=True)

    top = Inches(1.0)
    if subtitle:
        sbox = slide.shapes.add_textbox(CONTENT_X + Inches(0.12), Inches(0.86), CONTENT_W, Inches(0.36))
        sp = sbox.text_frame.paragraphs[0]
        sp.text = subtitle
        _para(sp, 10, GOLD)
        top = Inches(1.26)

    cols = len(headers)
    table_rows = 1 + len(rows)
    row_h = max(Inches(TABLE_ROW_MIN_IN), Inches(5.15 / max(table_rows, 1)))
    tw = CONTENT_W - Inches(0.12)
    table_shape = slide.shapes.add_table(table_rows, cols, CONTENT_X + Inches(0.08), top, tw, row_h * table_rows)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = GOLD_DIM
        _style_cell_paragraphs(cell, 9, CREAM, bold=True, align=PP_ALIGN.CENTER)

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL_ALT if r % 2 else NAVY_MID
            if c == len(row) - 1 and val.startswith("Rp"):
                _style_cell_paragraphs(cell, 8, GOLD, bold=True, align=PP_ALIGN.RIGHT)
            elif c == 0:
                _style_cell_paragraphs(cell, 8, CREAM)
            else:
                _style_cell_paragraphs(cell, 8, SILVER)

    add_footer(slide, prs)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- OPENING ---
    add_title_slide(
        prs,
        "INKAI Mobile Web",
        "Ekosistem Digital Institut Karate-Do Indonesia",
        "Platform: inkai-mobile-web.vercel.app\nVersi 0.1.0 | PWA Mobile-First | Next.js 16 + React 19",
        theme="hero",
    )

    add_bullet_slide(
        prs,
        "Agenda Presentasi",
        [
            "Latar Belakang & Permasalahan Organisasi",
            "Visi, Misi & Ringkasan Eksekutif",
            "Arsitektur Teknis & Struktur Organisasi Digital",
            "Modul Portal Publik (Beranda, Registrasi, Store)",
            "Modul Portal Anggota (Dashboard, Absensi, Iuran, Event, Prestasi)",
            "Modul Pertandingan Karate & UKT Digital (Roadmap)",
            "Modul Panel Administrator (Anggota, Verifikasi, Laporan)",
            "Fitur Unggulan, Manfaat Stakeholder & Roadmap",
            "Estimasi Biaya Development, Deploy & Operasional (www)",
            "Kesimpulan & Demo Platform",
        ],
        theme="agenda",
    )

    add_bullet_slide(
        prs,
        "Latar Belakang & Permasalahan",
        [
            "Organisasi karate nasional dengan ribuan anggota di puluhan cabang & ranting",
            "Pencatatan keanggotaan masih tersebar (Excel, WhatsApp, dokumen fisik)",
            "Verifikasi kenaikan tingkat, piagam & mutasi dojo memakan waktu lama",
            "Iuran bulanan dan biaya event sulit dipantau secara real-time",
            "Absensi latihan belum terdigitalisasi secara konsisten",
            "Komunikasi pengumuman organisasi tidak terpusat",
            "Anggota kesulitan mengakses informasi resmi & status keanggotaan",
        ],
        subtitle="Mengapa INKAI butuh transformasi digital?",
        theme="challenge",
    )

    add_bullet_slide(
        prs,
        "Solusi: INKAI Digital Ecosystem",
        [
            "Satu portal terpadu untuk anggota, pengurus, dan masyarakat",
            "Identitas digital resmi setiap anggota (NIA, kartu keanggotaan, profil lengkap)",
            "Administrasi otomatis dari tingkat ranting hingga pusat",
            "Transparansi iuran, absensi, dan verifikasi prestasi",
            "Komunikasi broadcast real-time pengurus → anggota",
            "Dapat diakses via browser — install sebagai PWA tanpa app store",
        ],
        subtitle="Satu Portal. Satu Identitas. Satu Organisasi.",
        theme="solution",
    )

    add_two_column_slide(
        prs,
        "Visi & Misi",
        "Visi",
        [
            "Menjadi platform digital terdepan untuk pengelolaan keanggotaan, kegiatan, dan administrasi INKAI di seluruh Indonesia",
        ],
        "Misi",
        [
            "Identitas digital resmi setiap anggota",
            "Sederhanakan administrasi organisasi",
            "Tingkatkan transparansi iuran & absensi",
            "Percepat komunikasi pengurus–anggota",
            "Dukung pertumbuhan berbasis data",
        ],
        theme="vision",
    )

    add_table_slide(
        prs,
        "Ringkasan Eksekutif — Tiga Lapisan Platform",
        ["Lapisan", "Pengguna", "Fungsi Utama"],
        [
            ["Portal Publik", "Masyarakat & calon anggota", "Informasi, katalog store, pendaftaran"],
            ["Portal Anggota", "Anggota terdaftar", "Dashboard, absensi, iuran, event, prestasi"],
            ["Panel Admin", "Pengurus & administrator", "Manajemen anggota, verifikasi, laporan"],
        ],
        theme="executive",
    )

    # --- TECH ---
    add_section_slide(prs, "BAGIAN I", "Arsitektur & Teknologi", "Infrastruktur modern, aman, dan scalable", theme="technology")

    add_table_slide(
        prs,
        "Stack Teknologi",
        ["Komponen", "Teknologi"],
        [
            ["Frontend", "Next.js 16, React 19, TypeScript, Framer Motion"],
            ["Styling", "CSS Modules, Glass Morphism, Tema Adaptif Siang/Malam"],
            ["Backend API", "REST API + JWT Authentication"],
            ["Storage", "Supabase (aset & dokumen)"],
            ["OCR", "Tesseract.js — baca kartu BPJS otomatis"],
            ["QR Scanner", "html5-qrcode — absensi real-time"],
            ["Deployment", "Vercel (CDN global, auto-scaling)"],
            ["PWA", "Manifest + Service Worker — installable"],
        ],
        theme="technology",
    )

    add_bullet_slide(
        prs,
        "Struktur Organisasi Digital",
        [
            "Pusat (Nasional) → Provinsi → Cabang → Dojo/Ranting → Anggota",
            "Administrator Pusat — akses penuh seluruh nasional",
            "Admin Provinsi — kelola cabang & ranting di provinsinya",
            "Admin Cabang — kelola dojo/ranting di cabangnya",
            "Admin Dojo/Ranting — kelola anggota di rantingnya",
            "Role-Based Access Control (RBAC) dengan scope data per wilayah",
        ],
        subtitle="Hierarki 5 level admin",
        theme="organization",
    )

    # --- PUBLIC ---
    add_section_slide(prs, "BAGIAN II", "Modul Portal Publik", "Akses tanpa login — informasi & onboarding", theme="public_web")

    add_bullet_slide(
        prs,
        "Beranda Publik (Landing Page)",
        [
            "Header sticky dengan branding INKAI Digital Ecosystem",
            "Carousel informasi otomatis dengan navigasi manual",
            "Tab navigasi dinamis — konten dari CMS backend",
            "Konten halaman publik berbasis Markdown",
            "Tombol Masuk / Dashboard (adaptif jika sudah login)",
            "Akses cepat ke INKAI Store resmi",
            "Tema adaptif siang/malam (Automatic Clock Theme)",
        ],
        subtitle="URL: /",
        theme="public_web",
    )

    add_bullet_slide(
        prs,
        "Registrasi & Autentikasi",
        [
            "Form pendaftaran: nama, email, HP, password",
            "Pemilihan lokasi berjenjang: Provinsi → Cabang → Dojo/Ranting",
            "Auto-login setelah registrasi berhasil",
            "Status awal PENDING — menunggu verifikasi admin & aktivasi NIA",
            "Login anggota dengan email/NIA + password",
            "Lupa password dengan token reset via email",
            "Panel admin terpisah di /admin/login",
        ],
        subtitle="URL: /register | /login | /forgot-password",
        theme="registration",
    )

    add_bullet_slide(
        prs,
        "INKAI Store (Publik)",
        [
            "Katalog produk resmi: peralatan & merchandise INKAI",
            "Dapat dijelajahi tanpa login",
            "Menampilkan nama, deskripsi, harga IDR, stok, dan foto",
            "Indikator stok habis",
            "Login diperlukan untuk checkout (fase berikutnya)",
            "Sinkron otomatis dengan modul inventaris admin",
        ],
        subtitle="URL: /store — Marketplace Resmi INKAI",
        theme="store",
    )

    # --- MEMBER ---
    add_section_slide(prs, "BAGIAN III", "Modul Portal Anggota", "Dashboard & layanan self-service anggota", theme="member")

    add_bullet_slide(
        prs,
        "Dashboard Anggota",
        [
            "Kartu Keanggotaan Digital — foto, NIA, tingkat sabuk, status",
            "Progress absensi bulan berjalan (X dari 8 sesi standar)",
            "Event terdekat & kegiatan saya yang sedang berjalan",
            "10 Quick Actions: Absensi, Iuran, Materi, Store, Sabuk, Piagam, Pelatihan, Pindah, Dokumen, Event",
            "Badge notifikasi belum dibaca",
            "Gate iuran: blokir event jika iuran bulanan belum lunas",
        ],
        subtitle="URL: /dashboard",
        theme="dashboard",
    )

    add_bullet_slide(
        prs,
        "Modul Absensi",
        [
            "Scanner QR Code real-time via kamera (html5-qrcode)",
            "Absensi dojo & event dengan validasi waktu otomatis",
            "Absensi event langsung dari daftar kegiatan saya",
            "Kalender absensi bulanan interaktif dengan indikator hari",
            "Riwayat absensi lengkap: waktu, metode, lokasi/dojo",
            "Aturan: maksimal 1 absensi per hari kalender",
        ],
        subtitle="URL: /absensi",
        theme="attendance",
    )

    add_bullet_slide(
        prs,
        "Modul Iuran & Pembayaran",
        [
            "Tagihan iuran bulanan (MONTHLY_IURAN) & biaya event (EVENT_FEE)",
            "Status: Belum Bayar, Menunggu Verifikasi, Lunas, Ditolak",
            "Filter tab, tahun, bulan & pencarian kata kunci",
            "Metode: Virtual Account, Transfer Bank, QR Code",
            "Upload bukti transfer dengan kompresi gambar otomatis",
            "Pembayaran multi-tagihan sekaligus",
            "Riwayat transaksi & preview bukti transfer",
        ],
        subtitle="URL: /billing",
        theme="billing",
    )

    add_bullet_slide(
        prs,
        "Modul Event & Kegiatan",
        [
            "Daftar event mendatang & riwayat event lalu",
            "Kategori otomatis: UKT/Ujian vs Turnamen/Kegiatan",
            "Detail event: deskripsi, jadwal, lokasi, cabang",
            "Pendaftaran online per kategori dengan tarif berbeda",
            "Countdown batas pendaftaran",
            "Pembayaran biaya event langsung di halaman detail",
            "Status: PENDING → DISETUJUI → LUNAS",
        ],
        subtitle="URL: /events | /events/[id]",
        theme="events",
    )

    # --- TOURNAMENT & UKT (Roadmap) ---
    add_section_slide(
        prs,
        "MODUL LANJUTAN",
        "Pertandingan Karate & UKT Digital",
        "Ekosistem turnamen real-time & ujian kenaikan tingkat terintegrasi — Fase 3 Roadmap",
        theme="tournament",
    )

    add_bullet_slide(
        prs,
        "Modul Pertandingan — Pendaftaran & Pengelompokan Kelas",
        [
            "Pendaftaran peserta online terhubung akun INKAI (NIA, dojo, tingkat Kyu/Dan)",
            "Validasi otomatis: usia, berat badan, gender, dan tingkat sabuk sesuai regulasi",
            "Pengelompokan otomatis ke kelas resmi (Kumite/Kata, Usia, BB, Kyu)",
            "Dukungan kelas Open — peserta lintas kategori dengan aturan fleksibel panitia",
            "Dukungan kelas Festival — pertandingan non-kompetitif / demonstrasi / hiburan",
            "Integrasi pembayaran biaya pendaftaran via modul billing INKAI",
            "Export daftar peserta per kelas untuk verifikasi wasit & medical check",
        ],
        subtitle="Modul Turnamen — Tahap Registrasi",
        theme="tournament",
    )

    add_bullet_slide(
        prs,
        "Modul Pertandingan — Drawing Bagan & Jadwal",
        [
            "Generate bagan (bracket) otomatis: single elimination, round-robin, pool-to-bracket",
            "Drawing acak atau seeded berdasarkan ranking/poin prestasi INKAI",
            "Penjadwalan pertandingan per kelas dengan estimasi durasi per babak",
            "Alokasi jadwal ke beberapa tatami secara paralel (Tatami 1, 2, 3, …)",
            "Durasi pertandingan configurable per kelas (Open, Festival, Kumite, Kata)",
            "Buffer waktu antar-pertandingan & jeda tatami untuk peralatan/medical",
            "Publish jadwal ke portal anggota, layar tatami, dan link publik live",
        ],
        subtitle="Modul Turnamen — Tahap Bagan & Scheduling",
        theme="tatami",
    )

    add_two_column_slide(
        prs,
        "Modul Pertandingan — Penilaian Juri & Live Real-time",
        "Penilaian Otomatis Juri",
        [
            "Interface scoring juri via tablet/mobile — login role Wasit/Juri",
            "Input poin Kumite (Yuko, Waza-ari, Ippon) & penalti otomatis",
            "Penilaian Kata: skor multi-juri, buang tertinggi/terendah, rata-rata",
            "Hasil babak langsung tersinkron ke bagan — pemenang maju otomatis",
            "Audit log setiap keputusan juri (transparansi & banding)",
        ],
        "Dashboard Real-time",
        [
            "Live board: pertandingan berjalan per tatami",
            "Countdown durasi pertandingan & sisa waktu per babak",
            "Antrean \"On Deck\" — peserta berikutnya per tatami",
            "Notifikasi push ke peserta saat giliran bertanding",
            "Leaderboard & medali per kelas Open / Festival / Resmi",
            "Update publik real-time tanpa refresh (WebSocket)",
        ],
        theme="scoring",
    )

    add_table_slide(
        prs,
        "Modul Pertandingan — Ringkasan Fitur per Kelas",
        ["Aspek", "Kelas Resmi", "Open / Festival"],
        [
            ["Pendaftaran", "Validasi strict Kyu/BB/Usia", "Aturan fleksibel panitia"],
            ["Bagan", "Bracket standar INKAI/WKF", "Round-robin / showcase"],
            ["Tatami", "Multi-tatami terjadwal", "Multi-tatami paralel"],
            ["Durasi", "Fixed per regulasi babak", "Configurable bebas"],
            ["Scoring", "Juri resmi + auto advance", "Demo / hiburan / non-poin"],
            ["Live", "Real-time scoreboard", "Real-time jadwal & lineup"],
        ],
        subtitle="Satu platform untuk turnamen kompetitif maupun festival karate",
        theme="tournament",
    )

    add_bullet_slide(
        prs,
        "UKT Digital — Ujian Kenaikan Tingkat Tanpa Raport",
        [
            "Ujian kenaikan tingkat (UKT) sepenuhnya digital — tanpa raport/ledger kertas",
            "Pendaftaran UKT terintegrasi modul event INKAI (bayar, absensi, verifikasi)",
            "Penilaian per materi ujian: kihon, kata, kumite, teori — input penguji via app",
            "Skor & keputusan Lulus/Tidak Lulus langsung masuk database anggota",
            "Kenaikan sabuk otomatis update profil, kartu keanggotaan, & riwayat rank",
            "Riwayat UKT tersimpan permanen — audit trail untuk cabang & pusat",
            "Sertifikat digital & notifikasi ke anggota setelah penguji submit hasil",
            "Tidak perlu re-entry manual — hilangnya duplikasi data Excel/WhatsApp",
        ],
        subtitle="Terintegrasi penuh dengan platform INKAI — Fase 3 Roadmap",
        theme="ukt",
    )

    add_two_column_slide(
        prs,
        "Integrasi UKT & Pertandingan ke Platform INKAI",
        "Data Terhubung Otomatis",
        [
            "Profil anggota (NIA, Kyu, dojo) — single source of truth",
            "Modul billing — biaya UKT & turnamen",
            "Modul absensi — kehadiran ujian & weigh-in",
            "Modul achievement — piagam & medali otomatis",
            "Modul verifikasi — approval kenaikan sabuk",
        ],
        "Manfaat Organisasi",
        [
            "Satu dashboard untuk UKT + turnamen + keanggotaan",
            "Laporan real-time: jumlah lulus, medali, peserta per cabang",
            "Standarisasi penilaian nasional tanpa raport fisik",
            "Transparansi hasil ujian & pertandingan ke anggota",
            "Data historis untuk seleksi atlet & UKT berikutnya",
        ],
        theme="ukt",
    )

    add_table_slide(
        prs,
        "Modul Prestasi & Achievement",
        ["Tab", "Jenis Pengajuan", "Contoh"],
        [
            ["Sabuk", "Kenaikan Tingkat (RANK_PROMOTION)", "Kyu 9 → Kyu 8"],
            ["Piagam", "Prestasi Pertandingan (PIAGAM)", "Juara 1 Kumite"],
            ["Pelatihan", "Sertifikasi (PELATIHAN)", "Kursus Wasit"],
        ],
        subtitle="URL: /achievement — Gate: wajib NIA & dokumen lengkap",
        theme="achievement",
    )

    add_bullet_slide(
        prs,
        "Modul Dokumen, Mutasi & Profil",
        [
            "Upload Akte Kelahiran — drag & drop, kompresi, auto-save",
            "Upload Kartu BPJS — OCR otomatis (nama, NIK, nomor kartu)",
            "Deteksi ketidaksesuaian data BPJS vs profil",
            "Pengajuan pindah dojo: Provinsi → Cabang → Dojo tujuan + alasan",
            "Riwayat mutasi & blokir pengajuan ganda",
            "Profil lengkap: NIK, BPJS, dojo, edit profil, ganti password",
            "Notifikasi in-app dengan mark-as-read",
        ],
        subtitle="URL: /documents | /transfer | /profile | /notifications",
        theme="documents",
    )

    add_bullet_slide(
        prs,
        "Modul Panduan & Materi",
        [
            "Panduan selamat datang anggota baru (konten dinamis admin)",
            "Overlay panduan otomatis saat pertama login",
            "Admin dapat edit langsung dari /admin/guide",
            "Materi Digital / Library — dalam pengembangan",
            "  Rencana: video tutorial, modul teori, silabus ujian per Kyu",
        ],
        subtitle="URL: /guide | /library",
        theme="guide",
    )

    # --- ADMIN ---
    add_section_slide(prs, "BAGIAN IV", "Modul Panel Administrator", "Manajemen organisasi multi-level", theme="admin")

    add_bullet_slide(
        prs,
        "Dashboard Admin",
        [
            "KPI: Total Anggota, Total Dojo/Cabang, Pending Verifikasi (scoped per role)",
            "Aksi cepat: Tambah Anggota, Laporan Absensi, Verifikasi Iuran",
            "Kelola Organisasi & Sebaran Anggota per Ranting",
            "Antrean prioritas: anggota pending NIA & klaim verifikasi",
            "Pencarian anggota terbaru by NIA/nama",
            "Pengumuman & broadcast langsung dari dashboard",
            "Pusat bantuan WhatsApp support INKAI",
        ],
        subtitle="URL: /admin",
        theme="admin",
    )

    add_bullet_slide(
        prs,
        "Manajemen Anggota",
        [
            "CRUD anggota — tambah, edit, hapus, detail lengkap",
            "Bulk import — unggah banyak anggota sekaligus",
            "Verifikasi pendaftaran & aktivasi NIA",
            "Provision login — buatkan akun untuk anggota",
            "Update tingkat sabuk (rank) resmi",
            "Pencarian & filter: NIA, nama, status, dojo",
            "Pengecualian iuran (allowEventWithoutDues)",
        ],
        subtitle="URL: /admin/members",
        theme="members",
    )

    add_bullet_slide(
        prs,
        "Organisasi & Verifikasi",
        [
            "Hierarki 4 level: Provinsi → Cabang → Dojo/Ranting — full CRUD",
            "Assign admin per unit organisasi",
            "Data bank rekening dojo, jadwal latihan, alamat, kontak",
            "Antrean Kerja — tab Dokumen: mutasi, kenaikan sabuk, piagam, pelatihan",
            "Antrean Kerja — tab Iuran: verifikasi bukti transfer",
            "Approve/Reject dengan catatan admin (audit trail)",
            "Preview bukti pendukung langsung di panel",
        ],
        subtitle="URL: /admin/organization | /admin/verification",
        theme="verification",
    )

    add_bullet_slide(
        prs,
        "Iuran, Event & Absensi Admin",
        [
            "Dashboard keuangan: total iuran masuk, pending, belum bayar",
            "Filter & pencarian tagihan by nama, NIA, dojo, status, tipe",
            "Verifikasi & hapus tagihan (koreksi)",
            "CRUD event: judul, tanggal, lokasi, kategori & tarif",
            "Daftar peserta, bulk register, update status pendaftaran",
            "Laporan absensi — filter tanggal, koreksi waktu, hapus catatan salah",
        ],
        subtitle="URL: /admin/billing | /admin/events | /admin/attendance",
        theme="billing",
    )

    add_bullet_slide(
        prs,
        "Broadcast, Statistik & Inventaris",
        [
            "Broadcast pengumuman resmi ke seluruh anggota (INFO/WARNING/SUCCESS)",
            "Notifikasi in-app real-time",
            "Sebaran anggota per ranting dengan breakdown tingkat Kyu",
            "Expand daftar anggota per Kyu per ranting",
            "Visualisasi warna sabuk (Putih, Kuning, Hijau, Biru, Cokelat, Hitam)",
            "CRUD inventaris produk store — sinkron ke katalog publik",
            "Manajemen role & permission (5 level admin)",
        ],
        subtitle="URL: /admin/broadcast | /admin/ranting-stats | /admin/inventory | /admin/settings",
        theme="broadcast",
    )

    add_bullet_slide(
        prs,
        "Admin: Kelola Pertandingan & UKT",
        [
            "Buat event turnamen / UKT dengan template kelas & tatami",
            "Kelola pendaftaran, approve peserta, assign ke kelas",
            "Trigger drawing bagan & publish jadwal multi-tatami",
            "Assign juri/wasit per tatami dengan hak akses scoring",
            "Monitor live scoreboard & override hasil (dengan audit log)",
            "Input & finalisasi hasil UKT — bulk update kenaikan sabuk",
            "Export laporan medali, statistik kelas Open/Festival, & rekap UKT",
        ],
        subtitle="URL: /admin/events (extended) — Panel Panitia Turnamen & Penguji",
        theme="tournament",
    )

    # --- HIGHLIGHTS ---
    add_section_slide(prs, "BAGIAN V", "Fitur Unggulan & Roadmap", "Differentiator & rencana pengembangan", theme="features")

    add_table_slide(
        prs,
        "12 Fitur Unggulan Platform",
        ["#", "Fitur", "Manfaat"],
        [
            ["1", "Kartu Keanggotaan Digital", "Identitas resmi INKAI di genggaman"],
            ["2", "QR Absensi", "Scan & absen dalam hitungan detik"],
            ["3", "OCR BPJS", "Baca otomatis data kartu BPJS"],
            ["4", "Multi-Level Admin", "Scope akses sesuai hierarki organisasi"],
            ["5", "Gate Iuran-Event", "Kedisiplinan iuran terintegrasi event"],
            ["6", "Verifikasi Terpusat", "Satu antrean untuk semua approval"],
            ["7", "PWA Installable", "Install tanpa app store"],
            ["8", "Tema Adaptif", "UI siang/malam otomatis"],
            ["9", "Glass Morphism UI", "Desain premium & elegan"],
            ["10", "Notifikasi Real-time", "Broadcast & alert instan"],
            ["11", "Turnamen Multi-Tatami", "Bagan, jadwal & skor juri live"],
            ["12", "UKT Digital Tanpa Raport", "Ujian kenaikan terintegrasi penuh"],
        ],
        theme="features",
    )

    add_bullet_slide(
        prs,
        "Alur Keanggotaan End-to-End",
        [
            "1. Daftar Online → Verifikasi Admin → Aktivasi NIA",
            "2. Upload Dokumen (Akte + BPJS/OCR)",
            "3. Absensi Latihan (QR) + Bayar Iuran Bulanan",
            "4. Daftar UKT Digital → Ujian → Kenaikan Sabuk otomatis (tanpa raport)",
            "5. Daftar Turnamen → Kelas/Drawing → Bertanding multi-tatami live",
            "6. Ajukan Piagam/Pelatihan → Admin Approve → Catat Prestasi",
            "7. (Opsional) Pindah Dojo → Admin Approve → Update Dojo",
        ],
        theme="workflow",
    )

    add_two_column_slide(
        prs,
        "Manfaat Bagi Stakeholder",
        "Bagi Anggota",
        [
            "Akses 24/7 status keanggotaan, iuran, event",
            "Pengajuan prestasi tanpa antre fisik",
            "Transparansi tagihan & riwayat pembayaran",
            "Absensi tanpa kertas",
        ],
        "Bagi Pengurus & Pusat",
        [
            "Dashboard KPI real-time per wilayah",
            "Verifikasi terpusat & efisien",
            "Laporan absensi & iuran otomatis",
            "Data demografi anggota per ranting/Kyu",
            "Standarisasi proses nasional & audit trail",
        ],
        theme="stakeholder",
    )

    add_table_slide(
        prs,
        "Roadmap Pengembangan",
        ["Fase", "Modul", "Status"],
        [
            ["Fase 1", "Auth, Dashboard, Profil, Dokumen, Absensi QR", "✅ Live"],
            ["Fase 1", "Iuran, Event, Verifikasi, Admin Panel", "✅ Live"],
            ["Fase 1", "Store Katalog, Inventaris, Broadcast", "✅ Live"],
            ["Fase 2", "Store Checkout & Payment Gateway", "🔄 Development"],
            ["Fase 2", "Materi Digital / Library", "🔄 Development"],
            ["Fase 2", "Chat Anggota Real-time", "🔄 Development"],
            ["Fase 3", "Modul Pertandingan: bagan, multi-tatami, scoring juri", "📋 Planned"],
            ["Fase 3", "UKT Digital tanpa raport — terintegrasi platform", "📋 Planned"],
            ["Fase 3", "Push Notification, Export PDF/Excel", "📋 Planned"],
            ["Fase 3", "Integrasi Midtrans/Xendit", "📋 Planned"],
            ["Fase 4", "Native App iOS/Android", "📋 Planned"],
        ],
        theme="roadmap",
    )

    add_bullet_slide(
        prs,
        "Keamanan & Desain UX",
        [
            "Autentikasi JWT dengan auto-expire & redirect",
            "RBAC 5 level dengan scope data per wilayah",
            "Validasi dokumen identitas (Akte + BPJS) + OCR cross-check",
            "Admin notes pada setiap approval (audit trail)",
            "HTTPS enforced via Vercel",
            "Mobile-first, dark premium theme, aksen emas INKAI",
            "Glass morphism, micro-animations, skeleton loading",
            "Bottom navigation ergonomis satu tangan",
        ],
        theme="security",
    )

    add_bullet_slide(
        prs,
        "Demo & Cara Akses",
        [
            "Buka: https://inkai-mobile-web.vercel.app/",
            "Publik: Jelajahi beranda, store, informasi INKAI",
            "Anggota: Klik Masuk → login akun terdaftar → /dashboard",
            "Admin: /admin/login → panel administrator",
            "Install PWA: Browser mobile → Add to Home Screen",
            "Halaman demo kunci: /absensi, /billing, /events, /admin",
        ],
        theme="demo",
    )

    # --- COST ESTIMATION ---
    add_section_slide(
        prs,
        "BAGIAN VI",
        "Estimasi Biaya Pengembangan & Operasional",
        "Rincian investasi development hingga go-live domain www — per Juni 2026",
        theme="cost",
    )

    add_bullet_slide(
        prs,
        "Kerangka Estimasi Biaya",
        [
            "Mata uang: Rupiah (IDR) — harga perkiraan pasar development web/mobile Indonesia",
            "Metode: estimasi berbasis modul aktual INKAI Mobile Web + Backend API",
            "Tim acuan: 1 Project Manager, 1 UI/UX, 2 Full-stack Developer, 1 QA",
            "Rate blended: Rp 850.000 – Rp 1.200.000 per hari kerja efektif",
            "Fase 1: sudah dibangun & live di inkai-mobile-web.vercel.app",
            "Fase 2–4: pengembangan lanjutan menuju produk lengkap + domain www resmi",
            "Biaya infrastruktur: recurring bulanan/tahunan (hosting, domain, database)",
            "Catatan: angka dapat disesuaikan scope, timeline, dan model kontrak (fixed/lump sum)",
        ],
        subtitle="Asumsi dasar perhitungan",
        theme="cost",
    )

    add_cost_table_slide(
        prs,
        "Biaya Development — Fase 1 (Sudah Live)",
        ["No", "Paket Pekerjaan / Modul", "Est. Hari", "Estimasi Biaya"],
        [
            ["1", "Analisis kebutuhan, wireframe, arsitektur API & database", "12", fmt_idr(14_400_000)],
            ["2", "UI/UX Design System (mobile-first, glass theme, komponen)", "18", fmt_idr(21_600_000)],
            ["3", "Backend API + Auth JWT + RBAC 5 level + Prisma/PostgreSQL", "45", fmt_idr(54_000_000)],
            ["4", "Portal Publik: landing, carousel, CMS tab, registrasi, login", "15", fmt_idr(18_000_000)],
            ["5", "Portal Anggota: dashboard, profil, dokumen, OCR BPJS", "28", fmt_idr(33_600_000)],
            ["6", "Modul Absensi QR + kalender + riwayat kehadiran", "14", fmt_idr(16_800_000)],
            ["7", "Modul Iuran & pembayaran (VA, transfer, bukti, multi-tagihan)", "20", fmt_idr(24_000_000)],
            ["8", "Modul Event & pendaftaran UKT/kegiatan + integrasi billing", "22", fmt_idr(26_400_000)],
            ["9", "Modul Achievement (sabuk, piagam, pelatihan) + mutasi dojo", "16", fmt_idr(19_200_000)],
            ["10", "Panel Admin: anggota, organisasi, verifikasi, billing, event", "40", fmt_idr(48_000_000)],
            ["11", "Panel Admin: absensi, broadcast, statistik ranting, inventaris", "18", fmt_idr(21_600_000)],
            ["12", "PWA (manifest, service worker), notifikasi, panduan anggota", "10", fmt_idr(12_000_000)],
            ["13", "QA, UAT, perbaikan bug, hardening keamanan", "20", fmt_idr(24_000_000)],
            ["14", "Dokumentasi teknis & pelatihan pengurus (online)", "8", fmt_idr(9_600_000)],
            ["", "SUBTOTAL FASE 1", "286 hari", fmt_idr(343_200_000)],
        ],
        subtitle="Status: ✅ Selesai & deployed (Vercel preview)",
        col_widths=[0.45, 7.2, 1.0, 1.8],
        theme="investment",
    )

    add_cost_table_slide(
        prs,
        "Biaya Development — Fase 2 s/d 4 (Pengembangan Lanjutan)",
        ["No", "Paket Pekerjaan / Modul", "Est. Hari", "Estimasi Biaya"],
        [
            ["1", "Store: checkout, keranjang, order management", "18", fmt_idr(21_600_000)],
            ["2", "Integrasi payment gateway (Midtrans / Xendit)", "14", fmt_idr(16_800_000)],
            ["3", "Materi Digital / Library (upload, kategori per Kyu)", "16", fmt_idr(19_200_000)],
            ["4", "Chat anggota real-time (WebSocket / Supabase Realtime)", "18", fmt_idr(21_600_000)],
            ["5", "Push notification (FCM / Web Push)", "10", fmt_idr(12_000_000)],
            ["6", "Export laporan PDF & Excel (absensi, iuran, anggota)", "12", fmt_idr(14_400_000)],
            ["7", "Modul Pertandingan Karate (registrasi, bagan, multi-tatami, scoring juri live)", "35", fmt_idr(42_000_000)],
            ["8", "UKT Digital tanpa raport (penilaian penguji, auto kenaikan sabuk)", "22", fmt_idr(26_400_000)],
            ["9", "Native app wrapper iOS & Android (Flutter/React Native)", "45", fmt_idr(54_000_000)],
            ["10", "QA regresi, UAT fase 2–4, optimasi performa", "15", fmt_idr(18_000_000)],
            ["", "SUBTOTAL FASE 2–4", "205 hari", fmt_idr(246_000_000)],
        ],
        subtitle="Status: 🔄 Sebagian berjalan | 📋 Sebagian direncanakan",
        col_widths=[0.45, 7.2, 1.0, 1.8],
        theme="investment",
    )

    add_cost_table_slide(
        prs,
        "Biaya Deployment & Go-Live Domain www",
        ["No", "Item Pekerjaan / Layanan", "Frekuensi", "Estimasi Biaya"],
        [
            ["1", "Registrasi domain resmi (contoh: inkai.or.id / inkai.co.id)", "Tahun ke-1", fmt_idr(350_000)],
            ["2", "Registrasi domain alternatif + proteksi brand (.com / .id)", "Opsional", fmt_idr(450_000)],
            ["3", "Konfigurasi DNS, subdomain (www, api, admin)", "Sekali", fmt_idr(3_500_000)],
            ["4", "Mapping custom domain Vercel (frontend + SSL otomatis)", "Sekali", fmt_idr(2_000_000)],
            ["5", "Deploy backend API production (Vercel / serverless)", "Sekali", fmt_idr(4_500_000)],
            ["6", "Setup environment production (env vars, secrets, CORS)", "Sekali", fmt_idr(2_500_000)],
            ["7", "Migrasi database production + backup strategy", "Sekali", fmt_idr(6_000_000)],
            ["8", "CI/CD pipeline (GitHub → auto deploy staging & prod)", "Sekali", fmt_idr(5_000_000)],
            ["9", "Pengetesan go-live, load test, smoke test www", "Sekali", fmt_idr(4_000_000)],
            ["10", "Hypercare / pendampingan 30 hari pasca go-live", "Sekali", fmt_idr(15_000_000)],
            ["", "SUBTOTAL DEPLOY & GO-LIVE www", "", fmt_idr(43_300_000)],
        ],
        subtitle="Target: www.[domain-inkai].id → production resmi (bukan subdomain Vercel saja)",
        col_widths=[0.45, 7.0, 1.5, 1.8],
        theme="deploy",
    )

    add_cost_table_slide(
        prs,
        "Biaya Infrastruktur & Operasional (Recurring)",
        ["No", "Layanan / Infrastruktur", "Paket", "Biaya / Bulan"],
        [
            ["1", "Vercel Pro — frontend INKAI Mobile Web (custom domain, bandwidth)", "Pro", fmt_idr(320_000)],
            ["2", "Vercel Pro / Serverless — backend API inkai-backend", "Pro", fmt_idr(320_000)],
            ["3", "Database PostgreSQL — Supabase Pro (storage, backup, auth)", "Pro", fmt_idr(400_000)],
            ["4", "Object storage dokumen & foto (Supabase Storage / CDN)", "Included+", fmt_idr(150_000)],
            ["5", "Email transaksional (reset password, notifikasi)", "Starter", fmt_idr(100_000)],
            ["6", "Monitoring & uptime (Sentry / Better Stack — opsional)", "Basic", fmt_idr(200_000)],
            ["7", "Domain renewal (amortisasi tahunan)", "—", fmt_idr(30_000)],
            ["8", "Payment gateway fee (Midtrans/Xendit — per transaksi)", "Variabel", "±2–3% / trx"],
            ["", "SUBTOTAL OPERASIONAL", "", fmt_idr(1_520_000)],
        ],
        subtitle="Estimasi bulanan setelah go-live www — belum termasuk fee transaksi payment gateway",
        col_widths=[0.45, 6.8, 1.4, 1.8],
        theme="infrastructure",
    )

    add_cost_table_slide(
        prs,
        "Biaya Pemeliharaan & Dukungan (Opsional)",
        ["No", "Layanan Dukungan", "Cakupan", "Estimasi Biaya"],
        [
            ["1", "Maintenance bulanan — bug fix & patch keamanan", "Retainer", fmt_idr(8_000_000)],
            ["2", "Minor enhancement (2–3 fitur kecil / bulan)", "Retainer", fmt_idr(12_000_000)],
            ["3", "Dukungan teknis pengurus (WhatsApp/ticket, jam kerja)", "Bulanan", fmt_idr(3_000_000)],
            ["4", "Backup audit & review keamanan tahunan", "Tahunan", fmt_idr(15_000_000)],
            ["5", "Pelatihan ulang admin cabang/ranting (batch baru)", "Per sesi", fmt_idr(5_000_000)],
            ["6", "Upgrade major version (Next.js, React, dependencies)", "Tahunan", fmt_idr(20_000_000)],
            ["", "PAKET MAINTENANCE STANDAR", "Per bulan", fmt_idr(23_000_000)],
        ],
        subtitle="Rekomendasi pasca go-live untuk menjaga platform tetap aman & up-to-date",
        col_widths=[0.45, 6.5, 1.6, 1.9],
        theme="maintenance",
    )

    add_cost_table_slide(
        prs,
        "Ringkasan Total Investasi",
        ["Komponen Biaya", "Keterangan", "Estimasi"],
        [
            ["Development Fase 1", "Modul inti — sudah live", fmt_idr(343_200_000)],
            ["Development Fase 2–4", "Fitur lanjutan, turnamen, UKT & native app", fmt_idr(246_000_000)],
            ["Deployment & Go-Live www", "Domain, DNS, CI/CD, hypercare", fmt_idr(43_300_000)],
            ["Infrastruktur (Tahun 1)", "12 × Rp 1,52 jt/bulan", fmt_idr(18_240_000)],
            ["", "TOTAL INVESTASI TAHUN PERTAMA", fmt_idr(650_740_000)],
            ["Maintenance (Opsional)", "Retainer standar × 12 bulan", fmt_idr(276_000_000)],
            ["", "GRAND TOTAL + MAINTENANCE", fmt_idr(926_740_000)],
        ],
        subtitle="Paket lengkap: build → deploy www → operasional 1 tahun + maintenance",
        col_widths=[3.2, 5.5, 2.0],
        theme="investment",
    )

    total_y1 = 650_740_000
    dp30 = int(total_y1 * 0.30)

    add_two_column_slide(
        prs,
        "Skema Pembayaran & Timeline (Usulan)",
        "Tahap Pembayaran Development",
        [
            f"DP 30% — kickoff & desain (±{fmt_idr(dp30)})",
            "30% — milestone Fase 1 selesai UAT",
            "25% — deploy www & go-live production",
            "15% — hypercare 30 hari selesai",
            "Maintenance: tagihan bulanan di awal bulan",
        ],
        "Timeline Perkiraan",
        [
            "Fase 1 (selesai): ~6–7 bulan dev + QA",
            "Deploy www & go-live: 2–3 minggu",
            "Fase 2 (checkout, library, chat): 3–4 bulan",
            "Fase 3 (push notif, export, payment): 2 bulan",
            "Fase 4 (native app): 3–4 bulan",
            "Total Fase 2–4: ±8–10 bulan paralel/sequential",
        ],
        theme="payment",
    )

    add_bullet_slide(
        prs,
        "Catatan Penting Estimasi Biaya",
        [
            "Angka di atas bersifat estimasi proposal; final quote mengikuti scope lock & kontrak",
            "Fase 1 senilai ±" + fmt_idr(343_200_000) + " sudah terealisasi dalam build saat ini",
            "Deploy www membutuhkan kepemilikan domain resmi INKAI (diserahkan pengurus pusat)",
            "Biaya payment gateway (Midtrans/Xendit) bersifat variabel per transaksi iuran/event",
            "Paket Vercel/Supabase free tier dapat dipakai untuk pilot; production www disarankan Pro",
            "Opsi hemat: tunda Fase 4 (native app) → hemat ±" + fmt_idr(54_000_000),
            "Opsi hemat: maintenance self-managed oleh tim internal INKAI",
        ],
        subtitle="Transparansi & fleksibilitas anggaran",
        theme="cost",
    )

    # --- CLOSING ---
    add_title_slide(
        prs,
        "Terima Kasih",
        "INKAI Digital Ecosystem",
        "Satu Portal. Satu Identitas. Satu Organisasi.\n\ninkai-mobile-web.vercel.app\nWhatsApp Support tersedia di Panel Admin",
        theme="closing",
    )

    return prs


def main() -> None:
    if not LOGO_PATH.is_file():
        raise FileNotFoundError(f"Logo not found: {LOGO_PATH}")
    import os

    refresh = os.environ.get("INKAI_PPT_REFRESH_SCREENSHOTS") == "1"
    if refresh or not asset_path("hero").is_file():
        print("Capturing slide visuals from inkai-mobile-web.vercel.app...")
        capture_all(force=refresh)
    else:
        print(f"Using cached screenshots in {asset_path('hero').parent}")
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    try:
        prs.save(str(OUTPUT))
        saved = OUTPUT
    except PermissionError:
        alt = OUTPUT.with_name(OUTPUT.stem + "-updated.pptx")
        prs.save(str(alt))
        saved = alt
        print(f"Note: {OUTPUT.name} sedang dibuka — disimpan ke {alt.name}")
    print(f"Saved: {saved}")
    print(f"Logo:  {LOGO_PATH}")
    print(f"Assets: {asset_path('hero').parent}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
